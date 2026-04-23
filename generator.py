# =============================================================================
# generator.py  —  SSPM PPT Slide Builder
# =============================================================================
# Slides produced:
#   1     — Summary Dashboard  (phase KPI strip + instances table + notes)
#   2     — Region-wise Summary
#   3…N   — Phase 1 Detail  (DYNAMIC: 14 apps/slide, auto-paginated)
#   N+1…M — Phase 2 Detail  (DYNAMIC: same)
#   M+1   — Blockers
#   M+2   — Milestones
# =============================================================================

import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import config
import data as manual_data

# ── Palette ───────────────────────────────────────────────────────────────────
def _rgb(k):
    r, g, b = config.THEME[k]; return RGBColor(r, g, b)

class C: pass
for _k in config.THEME: setattr(C, _k, _rgb(_k))

SLIDE_W, SLIDE_H  = 10.0, 5.625
CONTENT_BOTTOM    = 5.32   # usable height above footer
PHASE_PALETTE     = [C.accent, C.green, C.purple, C.teal, C.gray]

# All canonical statuses in display order
ALL_STATUSES = ["completed", "in_progress", "not_started", "future_request", "de_scoped"]


# ── Primitives ────────────────────────────────────────────────────────────────
def i(v): return Inches(v)

def add_rect(slide, x, y, w, h, fill, border=None, bw=0.75):
    s = slide.shapes.add_shape(1, i(x), i(y), i(w), i(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bw)
    else:      s.line.fill.background()
    return s

def add_text(slide, text, x, y, w, h, sz=9, bold=False, color=None,
             align=PP_ALIGN.LEFT, va="middle", italic=False, wrap=True):
    tb = slide.shapes.add_textbox(i(x), i(y), i(w), i(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.auto_size = None
    tf.vertical_anchor = {"middle": MSO_ANCHOR.MIDDLE, "top": MSO_ANCHOR.TOP,
                           "bottom": MSO_ANCHOR.BOTTOM}.get(va, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.name = "Calibri"; r.font.color.rgb = color or C.txt_dark

def add_line(slide, x1, y1, x2, y2, color=None, width=0.5):
    ln = slide.shapes.add_connector(1, i(x1), i(y1), i(x2), i(y2))
    ln.line.color.rgb = color or C.border; ln.line.width = Pt(width)

def pill(slide, x, y, w, h, label, color, sz=7.5):
    r, g, b = color[0], color[1], color[2]
    light = RGBColor(min(r+210,255), min(g+210,255), min(b+210,255))
    add_rect(slide, x, y, w, h, light, color, 0.75)
    add_text(slide, label, x, y, w, h, sz=sz, bold=True, color=color, align=PP_ALIGN.CENTER)

def mini_pill(slide, label, done, total, col, x, y, w=1.75):
    H = 0.44
    add_rect(slide, x, y, w, H, C.card, C.border, 0.75)
    add_text(slide, label, x+0.08, y+0.03, w-0.16, 0.16, sz=7, bold=True, color=C.txt_muted, va="top")
    add_text(slide, f"{done}/{total}", x+0.08, y+0.16, w-0.16, 0.20, sz=13, bold=True, color=col, va="middle")
    add_rect(slide, x+0.08, y+H-0.07, w-0.16, 0.04, C.border)
    if total > 0 and done > 0:
        add_rect(slide, x+0.08, y+H-0.07, max(0.05,(done/total)*(w-0.16)), 0.04, col)

def section_label(slide, text, y):
    add_text(slide, text, 0.18, y, 9.64, 0.18, sz=7.5, bold=True, color=C.txt_muted, va="middle")
    add_line(slide, 0.18, y+0.18, 9.82, y+0.18, C.border, 0.5)

def header(slide, subtitle, badge=None):
    add_rect(slide, 0, 0, SLIDE_W, 0.60, C.header)
    add_rect(slide, 0, 0, 0.18, 0.60, C.accent)
    add_text(slide, f"SSPM Initiative  |  {subtitle}",
             0.28, 0.06, 8.3, 0.48, sz=13, bold=True, color=C.white, va="middle")
    if badge:
        add_rect(slide, 8.7, 0.09, 1.10, 0.42, C.accent)
        add_text(slide, badge, 8.7, 0.09, 1.10, 0.42,
                 sz=9, bold=True, color=C.white, align=PP_ALIGN.CENTER)

def org_footer(slide):
    add_rect(slide, 0, 5.345, SLIDE_W, 0.28, C.card_alt)
    add_text(slide, "[ Add your org footer here ]",
             0.3, 5.35, 9.4, 0.25, sz=7, italic=True,
             color=C.border, align=PP_ALIGN.CENTER)

def _phase_color(phase_label):
    try:   return PHASE_PALETTE[config.PHASE_VALUES.index(phase_label) % len(PHASE_PALETTE)]
    except ValueError: return C.gray

def _phase_label(raw_phase):
    """Apply PHASE_DISPLAY_NAMES mapping; fall back to raw value."""
    return config.PHASE_DISPLAY_NAMES.get(raw_phase, raw_phase)

def _sc(status): return _rgb(config.STATUS_COLORS.get(status, "gray"))
def _sl(status): return config.STATUS_LABELS.get(status, status.replace("_"," ").title())


# ── Status bar helper (5 segments) ───────────────────────────────────────────
def _draw_status_bar(slide, x, y, w, h, counts, total):
    """Draw a proportional 5-segment colour bar from a status-count dict."""
    if total == 0: return
    cx = x
    for st in ALL_STATUSES:
        n = counts.get(st, 0)
        if n <= 0: continue
        seg_w = max(0.01, (n / total) * w)
        add_rect(slide, cx, y, seg_w, h, _sc(st))
        cx += seg_w


# ── Data helpers ──────────────────────────────────────────────────────────────

def _map_status_value(raw_val):
    """Map a raw JIRA dropdown string → canonical status key."""
    v = str(raw_val or "").lower().strip()
    for canonical, synonyms in config.STATUS_MAPPING.items():
        if any(s.lower() == v for s in synonyms):
            return canonical
    return "not_started"

def _display_status(canonical):
    """For detail-slide rows: completed+de_scoped → 'Completed', others keep label."""
    if canonical in config.COMPLETED_STATUSES:
        return "completed", "Completed"
    return canonical, _sl(canonical)

def _app_display_status(app_instances_in_phase):
    """Aggregate display status for one application's instances in a phase."""
    statuses = [inst.get("status","not_started") for inst in app_instances_in_phase]
    # If every instance is completed or de-scoped → Completed
    if all(s in config.COMPLETED_STATUSES for s in statuses):
        return "completed", "Completed"
    # If any is in-progress → In Progress
    if any(s == "in_progress" for s in statuses):
        return "in_progress", "Onboarding In Progress"
    if any(s == "future_request" for s in statuses):
        return "future_request", "Future Request"
    return "not_started", "Not Started"

def _build_phase_rows(applications, phase_summary):
    rows = []
    for ph in config.PHASE_VALUES:
        ph_apps  = [a for a in applications
                    if any(inst["phase"] == ph for inst in a["instances"])]

        if ph_apps:
            ph_insts = [inst for a in ph_apps for inst in a["instances"] if inst["phase"]==ph]
            st_counts = {s: sum(1 for inst in ph_insts if inst.get("status") == s)
                         for s in ALL_STATUSES}
            total_inst = len(ph_insts)
            rows.append({
                "phase":      ph,
                "total_apps": len(ph_apps),
                "prod_done":  sum(inst.get("prod_done",0)      for inst in ph_insts),
                "prod_pend":  sum(inst.get("prod_total",0) - inst.get("prod_done",0) for inst in ph_insts),
                "prod_total": sum(inst.get("prod_total",0)     for inst in ph_insts),
                "np_done":    sum(inst.get("non_prod_done",0)  for inst in ph_insts),
                "np_pend":    sum(inst.get("non_prod_total",0) - inst.get("non_prod_done",0) for inst in ph_insts),
                "np_total":   sum(inst.get("non_prod_total",0) for inst in ph_insts),
                "st_counts":  st_counts,
                "total_inst": total_inst,
                "in_scope":   ph in config.IN_SCOPE_PHASES,
            })
        elif ph in phase_summary:
            ps = phase_summary[ph]
            pd, pt = ps.get("prod_done",0), ps.get("prod_total",0)
            rows.append({
                "phase": ph, "total_apps": ps.get("total_apps",0),
                "prod_done": pd, "prod_pend": pt-pd, "prod_total": pt,
                "np_done": 0, "np_pend": 0, "np_total": 0,
                "st_counts": {"completed": pd, "in_progress":0,
                              "not_started": pt-pd, "future_request":0, "de_scoped":0},
                "total_inst": pt,
                "in_scope": ps.get("in_scope", False),
            })
    return rows

def _build_region_rows(applications):
    region_map = {}
    for a in applications:
        for inst in a["instances"]:
            r = inst.get("region","Unknown")
            m = region_map.setdefault(r, {"prod_done":0,"prod_pend":0,"prod_total":0,
                                          "np_done":0,"np_pend":0,"np_total":0})
            m["prod_done"]  += inst.get("prod_done",0)
            m["prod_pend"]  += inst.get("prod_total",0)  - inst.get("prod_done",0)
            m["prod_total"] += inst.get("prod_total",0)
            m["np_done"]    += inst.get("non_prod_done",0)
            m["np_pend"]    += inst.get("non_prod_total",0) - inst.get("non_prod_done",0)
            m["np_total"]   += inst.get("non_prod_total",0)
    ordered = [r for r in config.REGIONS if r in region_map]
    ordered += [r for r in region_map if r not in ordered]
    return [{"region": r, "region_lbl": r, **region_map[r]} for r in ordered]


# ── Instances table (shared by Slide 1 and Slide 2) ──────────────────────────

def _draw_instances_table(slide, rows, label_key, label_w, CY,
                           label_color_fn=None, has_np_fn=None):
    TX, TW  = 0.18, 9.64
    N_SUB   = 8          # 3 Prod + 3 NP + 2 Grand
    SUB_W   = (TW - label_w) / N_SUB
    GH, SH2, DR = 0.195, 0.155, 0.205

    grp_defs = [
        ("Prod Instances",     3, RGBColor(0x06,0x5a,0x8a)),
        ("Non-Prod Instances", 3, RGBColor(0x04,0x60,0x4a)),
        ("Grand Total",        2, RGBColor(0x2d,0x2d,0x6a)),
    ]
    add_rect(slide, TX, CY, label_w, GH+SH2, C.header)
    add_rect(slide, TX, CY, 0.04, GH+SH2, C.accent)
    add_text(slide, label_key.replace("_lbl","").capitalize(),
             TX+0.08, CY, label_w-0.10, GH+SH2, sz=8, bold=True, color=C.white, va="middle")
    gx = TX + label_w
    for glbl, ncols, gcol in grp_defs:
        gw = ncols * SUB_W
        add_rect(slide, gx, CY, gw-0.02, GH, gcol)
        add_text(slide, glbl, gx, CY, gw-0.02, GH,
                 sz=7.5, bold=True, color=C.white, align=PP_ALIGN.CENTER, va="middle")
        sub_lbls = ["Completed","Pending","Total"] if ncols==3 else ["Instances","Done"]
        for si, sl in enumerate(sub_lbls):
            SX = gx + si*SUB_W
            add_rect(slide, SX, CY+GH, SUB_W-0.01, SH2, C.card_alt, C.border, 0.4)
            add_text(slide, sl, SX, CY+GH, SUB_W-0.01, SH2,
                     sz=6, bold=True, color=C.txt_muted, align=PP_ALIGN.CENTER, va="middle")
        gx += gw

    for ri, row in enumerate(rows):
        RY  = CY + GH + SH2 + ri*DR
        bg  = C.card if ri%2==0 else C.card_alt
        tc  = label_color_fn(row) if label_color_fn else C.accent
        add_rect(slide, TX, RY, TW, DR, bg, C.border, 0.4)
        add_rect(slide, TX, RY, 0.04, DR, tc)
        add_text(slide, row[label_key], TX+0.08, RY, label_w-0.10, DR,
                 sz=8.5, bold=True, color=tc, va="middle")

        pd = row.get("prod_done",0); pp = row.get("prod_pend",0); pt = row.get("prod_total",0)
        nd = row.get("np_done",0);   np_= row.get("np_pend",0);  nt = row.get("np_total",0)
        has_np = (has_np_fn(row) if has_np_fn else nt>0)
        np_v   = [str(nd),str(np_),str(nt)] if has_np else ["N/A","N/A","N/A"]
        vals   = [str(pd),str(pp),str(pt)] + np_v + [str(pt+nt), str(pd+nd)]

        def _cell(val, x, is_done=False, is_na=False, is_total=False):
            fc = (C.gray if is_na else C.green if is_done and val not in ("0","N/A")
                  else C.txt_muted if val=="0" else C.header if is_total else C.txt_dark)
            add_text(slide, val, x, RY, SUB_W-0.01, DR, sz=8,
                     bold=(is_done or is_total) and not is_na, italic=is_na,
                     color=fc, align=PP_ALIGN.CENTER, va="middle")

        _cell(vals[0], TX+label_w+0*SUB_W, is_done=True)
        _cell(vals[1], TX+label_w+1*SUB_W)
        _cell(vals[2], TX+label_w+2*SUB_W, is_total=True)
        _cell(vals[3], TX+label_w+3*SUB_W, is_done=has_np, is_na=not has_np)
        _cell(vals[4], TX+label_w+4*SUB_W, is_na=not has_np)
        _cell(vals[5], TX+label_w+5*SUB_W, is_total=has_np, is_na=not has_np)
        _cell(vals[6], TX+label_w+6*SUB_W, is_total=True)
        _cell(vals[7], TX+label_w+7*SUB_W, is_done=True)

    # Total row
    TR_Y = CY + GH + SH2 + len(rows)*DR
    add_rect(slide, TX, TR_Y, TW, DR, C.header)
    add_rect(slide, TX, TR_Y, 0.04, DR, C.accent)
    add_text(slide, "TOTAL", TX+0.08, TR_Y, label_w-0.10, DR,
             sz=8.5, bold=True, color=C.accent, va="middle")
    tots = [sum(r.get(k,0) for r in rows)
            for k in ("prod_done","prod_pend","prod_total","np_done","np_pend","np_total")]
    tots += [tots[2]+tots[5], tots[0]+tots[3]]
    tr_clr = [RGBColor(0x34,0xd3,0x99),C.amber,C.white,
              RGBColor(0x34,0xd3,0x99),C.amber,C.white,
              C.white,RGBColor(0x34,0xd3,0x99)]
    for ci,(tv,tfc) in enumerate(zip(tots,tr_clr)):
        add_text(slide, str(tv), TX+label_w+ci*SUB_W, TR_Y, SUB_W-0.01, DR,
                 sz=8.5, bold=True, color=tfc, align=PP_ALIGN.CENTER, va="middle")
    return TR_Y + DR


# =============================================================================
# SLIDE 1 — Summary Dashboard
# =============================================================================
def build_slide1(prs, phase_rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C.bg
    header(slide, "SSPM Integration — Summary Dashboard")

    CY = 0.66
    section_label(slide, "UNIQUE APPLICATIONS — BY PHASE", CY)
    CY += 0.22

    # ── Phase KPI strip ───────────────────────────────────────────────────────
    STRIP_H = 1.20      # tall enough to fit 5-status breakdown without overlap
    TOTAL_W = 1.48
    GAP     = 0.05
    n       = len(phase_rows)
    TIER_W  = (9.64 - TOTAL_W - GAP*(n+1)) / max(n, 1)

    # Grand-total card
    grand_apps = sum(r["total_apps"]  for r in phase_rows)
    scope_apps = sum(r["total_apps"]  for r in phase_rows if r.get("in_scope"))
    grand_inst = sum(r["total_inst"]  for r in phase_rows)
    add_rect(slide, 0.18, CY, TOTAL_W, STRIP_H, C.header)
    add_rect(slide, 0.18, CY, 0.05, STRIP_H, C.accent)
    add_text(slide, "All Phases", 0.30, CY+0.04, TOTAL_W-0.14, 0.16,
             sz=7, bold=True, color=RGBColor(0x7d,0xd3,0xfc), va="top")
    add_text(slide, str(grand_apps), 0.30, CY+0.20, TOTAL_W-0.14, 0.32,
             sz=26, bold=True, color=C.white, va="middle")
    add_text(slide, "Applications", 0.30, CY+0.54, TOTAL_W-0.14, 0.14,
             sz=6, color=RGBColor(0x94,0xa3,0xb8), va="top")
    add_text(slide, f"{scope_apps} in-scope",
             0.30, CY+0.68, TOTAL_W-0.14, 0.12,
             sz=5.5, color=RGBColor(0x7d,0xd3,0xfc), va="top")
    add_text(slide, f"{grand_inst} instances",
             0.30, CY+0.80, TOTAL_W-0.14, 0.12,
             sz=5.5, color=RGBColor(0x7d,0xd3,0xfc), va="top")

    # Per-phase cards
    for ci, row in enumerate(phase_rows):
        X    = 0.18 + TOTAL_W + GAP + ci*(TIER_W + GAP)
        tc   = _phase_color(row["phase"])
        is_s = row.get("in_scope", False)
        total = row["total_apps"]
        st   = row["st_counts"]
        total_inst = row.get("total_inst", total)

        add_rect(slide, X, CY, TIER_W, STRIP_H, C.card, C.border, 0.75)
        add_rect(slide, X, CY, 0.05, STRIP_H, tc)

        # Phase badge — label from PHASE_BADGE_LABELS, fallback to IN/OUT SCOPE
        badge_lbl = config.PHASE_BADGE_LABELS.get(
            row["phase"], "IN SCOPE" if is_s else "OUT OF SCOPE"
        )
        bb  = RGBColor(0xdc,0xfc,0xe7) if is_s else RGBColor(0xf1,0xf5,0xf9)
        bfc = RGBColor(0x16,0x65,0x34) if is_s else RGBColor(0x47,0x55,0x69)
        add_rect(slide, X+TIER_W-1.04, CY+0.05, 1.00, 0.16, bb)
        add_text(slide, badge_lbl,
                 X+TIER_W-1.04, CY+0.05, 1.00, 0.16,
                 sz=5.5, bold=True, color=bfc, align=PP_ALIGN.CENTER, va="middle", wrap=False)

        ph_short = _phase_label(row["phase"]).replace("Phase ","P")
        add_text(slide, ph_short, X+0.10, CY+0.05, 0.55, 0.16,
                 sz=8, bold=True, color=tc, va="top")
        add_text(slide, str(total), X+0.10, CY+0.22, TIER_W-0.15, 0.30,
                 sz=24, bold=True, color=tc, va="middle")
        add_text(slide, "Unique Apps", X+0.10, CY+0.54, TIER_W-0.15, 0.13,
                 sz=6.5, color=C.txt_muted, va="top")

        # ── 5-status colour bar ───────────────────────────────────────────────
        BAR_Y = CY + 0.70
        BAR_H = 0.09
        bar_w = TIER_W - 0.16
        _draw_status_bar(slide, X+0.10, BAR_Y, bar_w, BAR_H, st, total_inst)

        # ── Status counts: count number + short label, one cell per status ────
        # Two rows: top=count (bold, larger), bottom=STATUS_SHORT (small)
        # wrap=False prevents text spilling outside the cell boundary
        parts = [(s, st.get(s, 0)) for s in ALL_STATUSES if st.get(s, 0) > 0]
        if parts:
            CELL_Y  = BAR_Y + BAR_H + 0.04
            CELL_H  = STRIP_H - (CELL_Y - CY) - 0.04   # fill remaining card height
            cell_w  = bar_w / len(parts)
            lx = X + 0.10
            for s_key, cnt in parts:
                short = config.STATUS_SHORT.get(s_key, s_key[:7])
                # Count number (top half)
                add_text(slide, str(cnt), lx, CELL_Y,
                         cell_w - 0.01, CELL_H * 0.55,
                         sz=9, bold=True, color=_sc(s_key),
                         align=PP_ALIGN.CENTER, va="middle", wrap=False)
                # Short label (bottom half)
                add_text(slide, short, lx, CELL_Y + CELL_H * 0.56,
                         cell_w - 0.01, CELL_H * 0.40,
                         sz=5, color=_sc(s_key),
                         align=PP_ALIGN.CENTER, va="top", wrap=False)
                lx += cell_w

    # ── Instances table ───────────────────────────────────────────────────────
    CY += STRIP_H + 0.10
    section_label(slide, "INSTANCES — PROD & NON-PROD ONBOARDING STATUS", CY)
    CY += 0.22

    for row in phase_rows:
        row["phase_lbl"] = _phase_label(row["phase"]).replace("Phase ","P")

    def ph_color(row): return _phase_color(row["phase"])
    def has_np(row):   return row.get("np_total",0) > 0

    bottom = _draw_instances_table(slide, phase_rows, "phase_lbl", 0.60, CY,
                                   label_color_fn=ph_color, has_np_fn=has_np)

    # ── Notes ─────────────────────────────────────────────────────────────────
    NB_Y = bottom + 0.06
    NB_H = 0.48
    NB_W = (9.64 - 0.10) / 2
    note_defs = [
        ("What is a Unique Application?",
         "A Unique Application is a distinct SaaS product in the SSPM programme "
         "(e.g. Salesforce, Workday). Counted once regardless of how many environments "
         "or regions it operates in.",
         "Example: Salesforce = 1 unique app, even with 8 regional instances.", C.accent),
        ("What is an Instance?",
         "An Instance is one JIRA Story representing a specific deployment — a regional "
         "environment (Prod or Non-Prod) or tenant. One app can have multiple instances, "
         "each onboarded into SSPM separately.",
         "Example: Salesforce has 8 instances — APAC Prod, EMEA Prod, NA Sandbox …", C.teal),
    ]
    for ni, (title, body, ex, nc) in enumerate(note_defs):
        NX = 0.18 + ni*(NB_W+0.10)
        add_rect(slide, NX, NB_Y, NB_W, NB_H, C.card, C.border, 0.75)
        add_rect(slide, NX, NB_Y, 0.04, NB_H, nc)
        add_text(slide, title, NX+0.10, NB_Y+0.04, NB_W-0.14, 0.16,
                 sz=7.5, bold=True, color=nc, va="top")
        add_text(slide, body,  NX+0.10, NB_Y+0.20, NB_W-0.14, 0.17,
                 sz=6.5, color=C.txt_mid, va="top")
        add_text(slide, ex,    NX+0.10, NB_Y+0.36, NB_W-0.14, 0.11,
                 sz=6, italic=True, color=C.txt_muted, va="top")

    # Status legend
    LEG_Y = NB_Y + NB_H + 0.04
    lx = 0.18
    for st_key in ALL_STATUSES:
        lbl = "● " + config.STATUS_SHORT.get(st_key, st_key)
        add_text(slide, lbl, lx, LEG_Y, 1.50, 0.15, sz=6.5, color=_sc(st_key))
        lx += 1.85

    org_footer(slide)


# =============================================================================
# SLIDE 2 — Region-wise Summary
# =============================================================================
def build_slide2_region(prs, region_rows):
    REG_COLORS = {"APAC":C.accent,"EMEA":C.teal,"Global":C.green,
                  "Japan":C.purple,"North America":C.amber,"TG":C.gray}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C.bg
    header(slide, "Integration Status — Region-wise Summary")
    CY = 0.68
    section_label(slide, "INSTANCE ONBOARDING STATUS BY REGION — PROD & NON-PROD", CY)
    CY += 0.22
    def rc(row): return REG_COLORS.get(row["region_lbl"], C.gray)
    def has_np(row): return row.get("np_total",0) > 0
    _draw_instances_table(slide, region_rows, "region_lbl", 1.20, CY,
                          label_color_fn=rc, has_np_fn=has_np)
    for lx, lbl, lc in [(0.18,"● Completed",C.green),(1.55,"● Pending",C.amber),
                         (2.85,"N/A — No Non-Prod environment",C.gray)]:
        add_text(slide, lbl, lx, 4.90, 3.5, 0.18, sz=6.5, color=lc)
    org_footer(slide)


# =============================================================================
# SLIDES 3…N — Phase Detail  (DYNAMIC — auto-paginated)
# =============================================================================

def build_detail_slides(prs, phase, applications):
    """
    Creates as many slides as needed for all apps in 'phase'.
    Automatically paginates at ROWS_PER_DETAIL_SLIDE rows.
    """
    # Build per-application rows for this phase
    app_rows = []
    for app in applications:
        ph_insts = [inst for inst in app["instances"] if inst.get("phase") == phase]
        if not ph_insts:
            continue
        disp_st, disp_lbl = _app_display_status(ph_insts)
        app_rows.append({
            "name":        app["name"],
            "disp_status": disp_st,
            "disp_label":  disp_lbl,
            "prod_done":   sum(inst.get("prod_done",0)      for inst in ph_insts),
            "prod_total":  sum(inst.get("prod_total",0)     for inst in ph_insts),
            "np_done":     sum(inst.get("non_prod_done",0)  for inst in ph_insts),
            "np_total":    sum(inst.get("non_prod_total",0) for inst in ph_insts),
            "n_instances": len(ph_insts),
        })

    if not app_rows:
        return

    RPP         = config.ROWS_PER_DETAIL_SLIDE
    total_pages = math.ceil(len(app_rows) / RPP)
    ph_color    = _phase_color(phase)
    ph_short    = _phase_label(phase).replace("Phase ","P")

    # Aggregate coverage for all apps in this phase (shown on every page)
    all_prod_done  = sum(r["prod_done"]  for r in app_rows)
    all_prod_total = sum(r["prod_total"] for r in app_rows)
    all_np_done    = sum(r["np_done"]    for r in app_rows)
    all_np_total   = sum(r["np_total"]   for r in app_rows)

    # 5-status counts at app level (for coverage bar on each page)
    all_st_counts = {s: sum(1 for r in app_rows
                            if _display_status(r["disp_status"])[0] == s
                            or r["disp_status"] == s)
                     for s in ALL_STATUSES}

    for page_idx in range(total_pages):
        page_rows = app_rows[page_idx*RPP : (page_idx+1)*RPP]
        _build_one_detail_slide(
            prs, phase, ph_short, ph_color,
            page_rows, page_idx+1, total_pages,
            len(app_rows),
            all_prod_done, all_prod_total,
            all_np_done,   all_np_total,
            all_st_counts,
        )


def _build_one_detail_slide(prs, phase, ph_short, ph_color,
                              page_rows, page_num, total_pages, total_apps,
                              all_prod_done, all_prod_total,
                              all_np_done, all_np_total, all_st_counts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C.bg

    badge = (f"{ph_short} · {total_apps} Apps"
             if total_pages == 1
             else f"{ph_short} · Pg {page_num}/{total_pages}")
    header(slide, f"{_phase_label(phase)} Applications — Integration Status", badge)

    # ── Coverage pills & 5-status bar ────────────────────────────────────────
    INST_Y = 0.66
    add_text(slide, "Instance Coverage",
             0.30, INST_Y+0.06, 1.55, 0.30, sz=8, bold=True, color=C.txt_muted, va="middle")
    mini_pill(slide, "PROD",     all_prod_done, all_prod_total, C.green, 2.00, INST_Y)
    if all_np_total > 0:
        mini_pill(slide, "NON-PROD", all_np_done,  all_np_total,  C.teal,  3.88, INST_Y)

    # 5-status bar top-right
    BAR_X, BAR_W = 6.00, 3.70
    add_text(slide, "App Status Breakdown",
             BAR_X, INST_Y+0.04, BAR_W, 0.16, sz=6.5, bold=True, color=C.txt_muted, va="top")
    _draw_status_bar(slide, BAR_X, INST_Y+0.20, BAR_W, 0.10,
                     all_st_counts, total_apps)
    # Status legend for the bar
    lx2 = BAR_X
    for st_key in ALL_STATUSES:
        n = all_st_counts.get(st_key, 0)
        if n > 0:
            short = config.STATUS_SHORT.get(st_key, st_key[:5])
            add_text(slide, f"{n} {short}", lx2, INST_Y+0.32, 0.80, 0.12,
                     sz=5.5, bold=True, color=_sc(st_key), va="top")
            lx2 += 0.76

    # ── Column layout ─────────────────────────────────────────────────────────
    # #(0.22) | App+3-boxes(3.30) | Status(2.10) | PROD(1.40) | NP(1.38) = 8.40 + margins
    C_NUM  = 0.30   # row number column x
    C_APP  = 0.55   # application name + 3-box area x
    C_STAT = 3.85   # status pill x
    C_PROD = 5.98   # prod pill x
    C_NP   = 7.42   # non-prod pill x
    W_APP  = 3.22   # app column width
    W_STAT = 2.05   # status column width
    W_PROD = 1.36   # prod column width
    W_NP   = 2.24   # np column width (to edge at 9.66)

    add_line(slide, 0.30, 1.18, 9.70, 1.18, C.border, 0.75)
    for cx, lbl, cw in [
        (C_NUM,  "#",                          0.22),
        (C_APP,  "Application  ·  Org / Prod / Non-Prod Instances", W_APP),
        (C_STAT, "Status",                     W_STAT),
        (C_PROD, "PROD",                       W_PROD),
        (C_NP,   "NON-PROD",                   W_NP),
    ]:
        add_text(slide, lbl, cx, 1.02, cw, 0.17,
                 sz=7, bold=True, color=C.txt_muted, align=PP_ALIGN.LEFT)

    # ── App rows ──────────────────────────────────────────────────────────────
    APP_Y  = 1.22
    NOTE_H = 0.32
    avail  = CONTENT_BOTTOM - NOTE_H - APP_Y
    ROW_H  = avail / config.ROWS_PER_DETAIL_SLIDE   # ~0.63 for 6 rows

    for idx, row in enumerate(page_rows):
        Y   = APP_Y + idx * ROW_H
        dst = row["disp_status"]
        sc  = _sc(dst)
        sl  = row["disp_label"]

        pd, pt = row["prod_done"],  row["prod_total"]
        nd, nt = row["np_done"],    row["np_total"]
        ni     = row["n_instances"]
        pc  = C.green if (pt > 0 and pd == pt) else (C.amber if pd > 0 else C.gray)
        nc  = C.teal  if (nt > 0 and nd == nt) else (C.amber if nd > 0 else C.gray)

        # Row background
        add_rect(slide, 0.30, Y, 9.40, ROW_H-0.03,
                 C.card if idx%2==0 else C.card_alt, C.border, 0.5)
        add_rect(slide, 0.30, Y, 0.05, ROW_H-0.03, sc)

        # Row number
        global_idx = (page_num-1)*config.ROWS_PER_DETAIL_SLIDE + idx + 1
        add_text(slide, str(global_idx), C_NUM+0.04, Y, 0.20, ROW_H-0.03,
                 sz=7, color=C.txt_muted, align=PP_ALIGN.CENTER, va="middle")

        # App name (upper ~48% of row)
        NAME_H = ROW_H * 0.46
        add_text(slide, row["name"], C_APP, Y + 0.05, W_APP - 0.06, NAME_H,
                 sz=9, bold=True, color=C.txt_dark, va="top")

        # ── 3 mini-boxes: ×orgs | ×Prod | ×NP — single line each, no wrap ──
        BOX_Y  = Y + NAME_H + 0.03
        BOX_H  = ROW_H - NAME_H - 0.10
        BOX_W  = (W_APP - 0.08) / 3
        box_defs = [
            (f"×{ni} Instance",  C.accent, RGBColor(0xDB, 0xEA, 0xFE)),
            (f"×{pt} Prod",      C.green,  RGBColor(0xD1, 0xFA, 0xE5)),
            (f"×{nt} Non-Prod",  C.teal,   RGBColor(0xCC, 0xFB, 0xF1)),
        ]
        bx = C_APP
        for blabel, bcol, blight in box_defs:
            add_rect(slide, bx, BOX_Y, BOX_W - 0.05, BOX_H, blight, bcol, 0.6)
            add_text(slide, blabel, bx + 0.02, BOX_Y, BOX_W - 0.09, BOX_H,
                     sz=7.5, bold=True, color=bcol,
                     align=PP_ALIGN.CENTER, va="middle", wrap=False)
            bx += BOX_W

        # Status pill  (binary: Completed/De-Scoped → green, others own colour)
        pill_h = max(0.22, ROW_H * 0.52)
        pill_y = Y + (ROW_H - 0.03 - pill_h) / 2
        pill(slide, C_STAT, pill_y, W_STAT - 0.06, pill_h, sl, sc, sz=7)

        # PROD pill
        if pt > 0:
            pill(slide, C_PROD, pill_y, W_PROD - 0.04, pill_h,
                 f"PROD {pd}/{pt}", pc, sz=7)
        else:
            add_text(slide, "—", C_PROD, Y, W_PROD - 0.04, ROW_H-0.03,
                     sz=9, color=C.txt_muted, align=PP_ALIGN.CENTER)

        # NP pill
        if nt > 0:
            pill(slide, C_NP, pill_y, W_NP - 0.04, pill_h,
                 f"NP {nd}/{nt}", nc, sz=7)
        else:
            add_text(slide, "—", C_NP, Y, W_NP - 0.04, ROW_H-0.03,
                     sz=9, color=C.txt_muted, align=PP_ALIGN.CENTER)

    # ── Note at bottom ────────────────────────────────────────────────────────
    NOTE_Y = CONTENT_BOTTOM - NOTE_H
    add_rect(slide, 0.18, NOTE_Y, 9.64, NOTE_H, RGBColor(0xFE,0xF9,0xC3),
             RGBColor(0xCA,0x8A,0x04), 0.75)
    add_rect(slide, 0.18, NOTE_Y, 0.05, NOTE_H, RGBColor(0xCA,0x8A,0x04))
    add_text(slide, "NOTE",
             0.30, NOTE_Y+0.03, 0.65, 0.14, sz=7, bold=True,
             color=RGBColor(0x92,0x40,0x0E), va="top")
    note_body = (
        "Status Grouping:  "
        "Completed  =  'Completed' + 'De-Scoped'  (shown in green).   "
        "Pending  =  'Onboarding In Progress' + 'Not Started' + 'Future Request'  "
        "(shown with their respective colour).   "
        "Status is determined by the JIRA dropdown field configured in config.py → STATUS_FIELD."
    )
    add_text(slide, note_body,
             0.30, NOTE_Y+0.05, 9.30, NOTE_H-0.08, sz=6.5,
             color=RGBColor(0x78,0x35,0x0F), va="middle")

    org_footer(slide)


# =============================================================================
# BLOCKERS SLIDE
# =============================================================================
def build_slide_blockers(prs, blockers):
    PHASE_COL = {ph: _phase_color(ph) for ph in config.PHASE_VALUES}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C.bg
    header(slide, "Blockers & Impediments", f"{len(blockers)} Blockers")
    CY = 0.68
    add_rect(slide, 0.18, CY, 9.64, 0.26, RGBColor(0xFF,0xED,0xED), C.red, 0.75)
    add_rect(slide, 0.18, CY, 0.06, 0.26, C.red)
    add_text(slide, "Active Blockers & Impediments",
             0.30, CY, 7.0, 0.26, sz=9.5, bold=True, color=C.red, va="middle")
    CY += 0.28
    BRH = 0.225
    for b in blockers:
        if CY + BRH > CONTENT_BOTTOM: break
        ic = C.red if b.get("impact","Med")=="High" else C.amber
        pc = PHASE_COL.get(b.get("phase",""), C.gray)
        add_rect(slide, 0.18, CY, 9.64, BRH-0.02, C.card, C.border, 0.5)
        add_rect(slide, 0.18, CY, 0.06, BRH-0.02, ic)
        add_text(slide, b.get("id",""),   0.30, CY, 0.40, BRH-0.02, sz=8, bold=True, color=ic, va="middle")
        add_text(slide, b.get("text",""), 0.74, CY, 5.60, BRH-0.02, sz=8.5, color=C.txt_dark, va="middle")
        ph = b.get("phase","")
        if ph: pill(slide, 6.40, CY+0.03, 1.10, BRH-0.08, ph, pc)
        add_text(slide, "Owner: "+b.get("owner","TBD"),
                 7.60, CY, 1.40, BRH-0.02, sz=7.5, color=C.txt_muted, align=PP_ALIGN.RIGHT, va="middle")
        pill(slide, 9.16, CY+0.03, 0.58, BRH-0.08, b.get("impact","Med"), ic)
        CY += BRH
    org_footer(slide)


# =============================================================================
# MILESTONES SLIDE  (auto-paginated — handles any number of tasks)
# =============================================================================
def build_slide_milestones(prs, milestones):
    TX, TW = 0.18, 9.64
    THdr, TRH = 0.28, 0.255

    def _new_ms_slide(page_num):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        sl.background.fill.solid(); sl.background.fill.fore_color.rgb = C.bg
        badge = f"Pg {page_num}" if page_num > 1 else None
        header(sl, "Milestones & Roadmap", badge)
        return sl, 0.68

    slide, CY = _new_ms_slide(1)
    page = 1

    for m in milestones:
        mc   = _phase_color(m.get("phase", "Phase 1"))
        sc   = _sc(m.get("status", "not_started"))
        r, g, b = mc[0], mc[1], mc[2]
        lbg  = RGBColor(min(r+215,255), min(g+215,255), min(b+215,255))
        tasks = m.get("tasks", [])

        # Start new slide if milestone header won't fit
        if CY + THdr > CONTENT_BOTTOM:
            org_footer(slide)
            page += 1
            slide, CY = _new_ms_slide(page)

        # Milestone section header
        add_rect(slide, TX, CY, TW, THdr, lbg, mc, 0.75)
        add_rect(slide, TX, CY, 0.06, THdr, mc)
        add_text(slide, m.get("name",""), TX+0.14, CY, 7.0, THdr,
                 sz=9.5, bold=True, color=mc, va="middle")
        pill(slide, TX+TW-1.55, CY+0.04, 1.45, THdr-0.08,
             _sl(m.get("status","not_started")), sc)
        CY += THdr

        # Task rows — overflow to new slide when needed
        for ri, task in enumerate(tasks):
            if CY + TRH > CONTENT_BOTTOM:
                org_footer(slide)
                page += 1
                slide, CY = _new_ms_slide(page)
                # Continuation header (compact)
                add_rect(slide, TX, CY, TW, THdr*0.65, lbg, mc, 0.75)
                add_text(slide, m.get("name","") + "  (continued)",
                         TX+0.14, CY, 8.0, THdr*0.65,
                         sz=8, bold=True, color=mc, va="middle")
                CY += THdr * 0.65

            ts = _sc(task.get("status", "not_started"))
            add_rect(slide, TX, CY, TW, TRH-0.03,
                     C.card if ri%2==0 else C.card_alt, C.border, 0.5)
            add_text(slide, str(ri+1), TX+0.10, CY, 0.28, TRH-0.03,
                     sz=8, bold=True, color=mc,
                     align=PP_ALIGN.CENTER, va="middle")
            add_text(slide, task.get("task",""), TX+0.42, CY, 7.1, TRH-0.03,
                     sz=9, color=C.txt_dark, va="middle")
            pill(slide, TX+TW-1.55, CY+0.04, 1.45, TRH-0.11,
                 _sl(task.get("status","not_started")), ts)
            CY += TRH

        CY += 0.08

    org_footer(slide)


# =============================================================================
# APPENDIX — Detailed instance / sub-task table  (auto-paginated)
# =============================================================================

def _get_appendix_rows(applications):
    """Flatten all instances + their sub-task rows into one list for the appendix."""
    rows = []
    for app in applications:
        for inst in app["instances"]:
            base = {
                "app_name"  : app["name"],
                "inst_name" : inst.get("instance_name") or inst.get("summary", ""),
                "app_id"    : inst.get("app_id", "—"),
                "phase"     : _phase_label(inst.get("phase", "")),
            }
            subtasks = inst.get("subtasks", [])
            if subtasks:
                for st in subtasks:
                    raw_st = st.get("status", "")
                    # Map raw sub-task status string → canonical key
                    raw_lower = raw_st.lower().strip()
                    canonical = "not_started"
                    for key, synonyms in config.STATUS_MAPPING.items():
                        if any(s.lower() == raw_lower for s in synonyms):
                            canonical = key
                            break
                    if st.get("done"):
                        canonical = "completed"
                    rows.append({**base,
                        "sub_key": st.get("key", ""),
                        "sub_name": st.get("name", ""),
                        "env"    : st.get("env", "—"),
                        "status" : canonical,
                    })
            else:
                # No sub-task detail — show one summary row per instance
                rows.append({**base,
                    "sub_key" : inst.get("key", ""),
                    "sub_name": inst.get("instance_name") or inst.get("summary", ""),
                    "env"     : "Prod + Non-Prod",
                    "status"  : inst.get("status", "not_started"),
                })
    return rows


def build_appendix_slides(prs, appendix_rows):
    rows = appendix_rows
    if not rows:
        return

    # Column definitions: (header, width_inches)
    COLS = [
        ("#",            0.30),
        ("Application",  1.76),
        ("Instance",     2.10),
        ("App ID",       0.90),
        ("Phase",        0.82),
        ("Env",          0.82),
        ("Sub-task Key", 0.94),
        ("Status",       1.64),
    ]
    TX     = 0.18
    HDR_H  = 0.26
    ROW_H  = 0.21
    START_Y = 0.68 + HDR_H
    RPP    = int((CONTENT_BOTTOM - START_Y) / ROW_H)   # rows per page
    total_pages = math.ceil(len(rows) / RPP)

    for page in range(total_pages):
        page_rows = rows[page * RPP : (page + 1) * RPP]

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = C.bg

        badge = f"Pg {page+1}/{total_pages}"
        header(slide, "Appendix — Application Instance Detail", badge)

        # ── Column headers ────────────────────────────────────────────────────
        CX = TX
        for lbl, cw in COLS:
            add_rect(slide, CX, 0.68, cw - 0.02, HDR_H, C.header)
            add_rect(slide, CX, 0.68, 0.03, HDR_H, C.accent)
            add_text(slide, lbl, CX + 0.06, 0.68, cw - 0.08, HDR_H,
                     sz=7, bold=True, color=C.white, va="middle", wrap=False)
            CX += cw

        # ── Data rows ─────────────────────────────────────────────────────────
        prev_app = None
        for ri, row in enumerate(page_rows):
            RY = START_Y + ri * ROW_H
            # Alternate background; shade differently when app group changes
            same_app = row["app_name"] == prev_app
            bg = C.card if same_app else RGBColor(0xF0, 0xF4, 0xFF) if ri % 2 == 0 else C.card_alt
            prev_app = row["app_name"]

            st_key = row["status"]
            sc     = _sc(st_key)
            sl     = config.STATUS_LABELS.get(st_key, st_key)

            add_rect(slide, TX, RY,
                     sum(cw for _, cw in COLS) - 0.02, ROW_H - 0.02,
                     bg, C.border, 0.3)

            global_row = page * RPP + ri + 1
            CX = TX
            for ci, (lbl, cw) in enumerate(COLS):
                val = [
                    str(global_row),
                    row["app_name"],
                    row["inst_name"],
                    row["app_id"],
                    row["phase"],
                    row["env"],
                    row["sub_key"],
                    None,            # status — drawn as pill
                ][ci]

                if lbl == "Status":
                    pill(slide, CX + 0.02, RY + 0.02,
                         cw - 0.06, ROW_H - 0.05, sl, sc, sz=6)
                elif lbl == "#":
                    add_text(slide, val, CX, RY, cw - 0.02, ROW_H,
                             sz=6.5, color=C.txt_muted,
                             align=PP_ALIGN.CENTER, va="middle", wrap=False)
                else:
                    add_text(slide, val, CX + 0.04, RY, cw - 0.06, ROW_H,
                             sz=7, color=C.txt_dark, va="middle", wrap=False)
                CX += cw

        org_footer(slide)


# =============================================================================
# MAIN
# =============================================================================
def main():
    if config.USE_JIRA:
        from jira_loader import load_from_jira
        d             = load_from_jira()
        slide1_apps   = d["slide1_apps"]    # from Stories  → Slide 1 KPI + region
        slide3_apps   = d["slide3_apps"]    # from Sub-tasks → Slide 3 detail
        appendix_rows = d["appendix_rows"]  # from Sub-tasks → Appendix flat list
        phase_summary = {}
        extra         = d["extra_sections"]
        raw_blockers  = extra.get("blockers", [])
        if raw_blockers:
            bf = config.EXTRA_SECTIONS.get("blockers",{}).get("fields",{})
            blockers = [{
                "id"    : f"B{i+1}",
                "key"   : iss.get("key",""),
                "text"  : (iss.get("fields") or {}).get("summary",""),
                "owner" : str(((iss.get("fields") or {}).get("assignee") or {}).get("displayName","TBD")),
                "impact": str(((iss.get("fields") or {}).get(bf.get("impact_field","")) or {}).get("value","Med")),
                "phase" : str(((iss.get("fields") or {}).get(bf.get("phase_field",""))  or {}).get("value","")),
            } for i,iss in enumerate(raw_blockers)]
        else:
            blockers = manual_data.BLOCKERS
        raw_milestones = extra.get("milestones", [])
        milestones = raw_milestones if raw_milestones else manual_data.MILESTONES
        print("  Data source: JIRA")
    else:
        # Manual mode: use APPLICATIONS for slide1, SLIDE3_APPLICATIONS for slide3
        slide1_apps   = manual_data.APPLICATIONS
        slide3_apps   = getattr(manual_data, "SLIDE3_APPLICATIONS", manual_data.APPLICATIONS)
        appendix_rows = manual_data.APPENDIX_ROWS if hasattr(manual_data, "APPENDIX_ROWS") else []
        phase_summary = manual_data.PHASE_SUMMARY
        blockers      = manual_data.BLOCKERS
        milestones    = manual_data.MILESTONES
        print("  Data source: data.py (manual)")

    # Slide 1 uses slide1_apps (stories); Slide 3 uses slide3_apps (sub-tasks)
    phase_rows  = _build_phase_rows(slide1_apps, phase_summary)
    region_rows = _build_region_rows(slide1_apps)

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    print("  Slide 1 — Summary Dashboard ...")
    build_slide1(prs, phase_rows)

    print("  Slide 2 — Region Summary ...")
    build_slide2_region(prs, region_rows)

    slide_num = 3
    for phase in config.IN_SCOPE_PHASES:
        ph_apps = [a for a in slide3_apps
                   if any(inst["phase"]==phase for inst in a["instances"])]
        if not ph_apps: continue
        n_pages = math.ceil(len(ph_apps) / config.ROWS_PER_DETAIL_SLIDE)
        print(f"  Slides {slide_num}–{slide_num+n_pages-1} — {phase} Detail "
              f"({len(ph_apps)} apps → {n_pages} slide{'s' if n_pages>1 else ''}) ...")
        build_detail_slides(prs, phase, slide3_apps)
        slide_num += n_pages

    print(f"  Slide {slide_num} — Blockers ...")
    build_slide_blockers(prs, blockers)
    slide_num += 1

    ms_start = slide_num
    build_slide_milestones(prs, milestones)
    ms_pages = len(prs.slides) - (ms_start - 1)
    if ms_pages > 1:
        print(f"  Slides {ms_start}–{ms_start+ms_pages-1} — Milestones ({ms_pages} slides, auto-paginated) ...")
    else:
        print(f"  Slide {ms_start} — Milestones ...")

    ax_start = len(prs.slides) + 1
    build_appendix_slides(prs, appendix_rows)
    ax_pages = len(prs.slides) - ax_start + 1
    if ax_pages > 0:
        print(f"  Slides {ax_start}–{ax_start+ax_pages-1} — Appendix ({ax_pages} slides, auto-paginated) ...")

    prs.save(config.OUTPUT_FILE)
    print(f"\nSaved → {config.OUTPUT_FILE}  ({len(prs.slides)} slides total)")


if __name__ == "__main__":
    print("\nSSPM PPT Generator")
    print("=" * 40)
    main()
